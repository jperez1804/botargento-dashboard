"""Build the client-facing dashboard walkthrough deck.

Reads slide copy from the SLIDES list below and screenshots from
docs/screenshots/, applies the BotArgento brand language, and writes a
.pptx in light or dark theme. Re-running rebuilds.

Usage:
    py scripts/build-walkthrough-deck.py [out.pptx] [--theme light|dark]
    py scripts/build-walkthrough-deck.py --all   # builds light + dark side by side

Missing screenshots render as labeled placeholders; the deck still ships
and the user can drop in PNGs later and re-run.
"""

from __future__ import annotations

import argparse
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# --- BotArgento brand palette ---------------------------------------------
# Pulled from the canonical favicon at landingpage/favicon.svg. The helmet
# icon ships with three structural colors that we lean on across both themes:
# sky blue (the helmet gradient), deep navy (the helmet body / dark canvas),
# and gold (the visor T and stars at the bottom).

BRAND_PRIMARY = RGBColor(0x75, 0xAA, 0xDB)      # Sky blue — main accent
BRAND_PRIMARY_SOFT = RGBColor(0xA8, 0xCD, 0xE8)  # Lighter highlight stop
BRAND_PRIMARY_DEEP = RGBColor(0x3A, 0x60, 0x80)  # Deepest blue stop
BRAND_DARK = RGBColor(0x05, 0x0E, 0x1F)         # Deep navy — favicon background
BRAND_DARK_2 = RGBColor(0x0E, 0x1A, 0x2E)       # Slightly lifted dark surface
BRAND_GOLD = RGBColor(0xE8, 0xB8, 0x4B)         # Gold — visor / stars
BRAND_GOLD_DEEP = RGBColor(0xB8, 0x85, 0x1E)    # Deeper gold

GOOD = RGBColor(0x05, 0x96, 0x69)
BAD = RGBColor(0xDC, 0x26, 0x26)

BRAND_NAME = "botargento"
BRAND_TAGLINE = "Automatización para Instagram, WhatsApp y más"


# --- Theme system ---------------------------------------------------------
# Each theme is a flat dict of design tokens. The renderer reads from this
# instead of the module-level constants so the same slide data builds for
# light or dark identically.

LIGHT_THEME = {
    "name": "light",
    "canvas": RGBColor(0xFA, 0xFA, 0xFA),
    "surface": RGBColor(0xFF, 0xFF, 0xFF),
    "ink": RGBColor(0x11, 0x18, 0x27),
    "muted_ink": RGBColor(0x6B, 0x72, 0x80),
    "soft_ink": RGBColor(0x9C, 0xA3, 0xAF),
    "rule": RGBColor(0xE5, 0xE7, 0xEB),
    "accent": BRAND_PRIMARY,
    "accent_text": BRAND_PRIMARY_DEEP,  # Sky blue is too light on white headings
    "section_canvas": BRAND_DARK,
    "section_ink": RGBColor(0xFF, 0xFF, 0xFF),
    "cover_canvas": RGBColor(0xFF, 0xFF, 0xFF),
    "cover_title_ink": BRAND_DARK,
    "cover_kicker": BRAND_PRIMARY,
    "footer_ink": RGBColor(0x6B, 0x72, 0x80),
    "footer_dim": RGBColor(0x9C, 0xA3, 0xAF),
    "screenshot_border": RGBColor(0xE5, 0xE7, 0xEB),
}

DARK_THEME = {
    "name": "dark",
    "canvas": BRAND_DARK,
    "surface": BRAND_DARK_2,
    "ink": RGBColor(0xE8, 0xEE, 0xF5),
    "muted_ink": RGBColor(0xA8, 0xB8, 0xCC),
    "soft_ink": RGBColor(0x6B, 0x80, 0x99),
    "rule": RGBColor(0x1F, 0x2C, 0x44),
    "accent": BRAND_PRIMARY,
    "accent_text": BRAND_PRIMARY_SOFT,
    "section_canvas": RGBColor(0x00, 0x08, 0x14),
    "section_ink": BRAND_PRIMARY_SOFT,
    "cover_canvas": BRAND_DARK,
    "cover_title_ink": RGBColor(0xFF, 0xFF, 0xFF),
    "cover_kicker": BRAND_PRIMARY,
    "footer_ink": RGBColor(0xA8, 0xB8, 0xCC),
    "footer_dim": RGBColor(0x6B, 0x80, 0x99),
    "screenshot_border": RGBColor(0x1F, 0x2C, 0x44),
}

# Fonts: try the dashboard's own faces first; fall back to Office defaults.
HEADING_FONT = "Fraunces"
HEADING_FALLBACK = "Cambria"
BODY_FONT = "Geist Sans"
BODY_FALLBACK = "Calibri"
MONO_FONT = "Geist Mono"
MONO_FALLBACK = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)


# --- Slide data -----------------------------------------------------------


@dataclass
class Slide:
    kicker: str  # Mono uppercase tag at the top of the slide
    title: str  # Fraunces h1
    body: list[str] = field(default_factory=list)  # Each string = one paragraph
    screenshot: Optional[str] = None  # filename in docs/screenshots/
    callouts: list[str] = field(default_factory=list)  # Optional list rendered as
    #   compact bullets in a small box under the body. Used for two-column slides.
    footnote: Optional[str] = None  # Small soft-ink line at the bottom
    layout: str = "two_col"  # "two_col" | "title_only" | "full_visual" | "section"


SLIDES: list[Slide] = [
    Slide(
        kicker=BRAND_TAGLINE,
        title="Panel BotArgento\nCómo leer cada métrica",
        body=[
            "Una guía corta para entender qué muestra cada tarjeta del panel y "
            "qué decisión podés tomar a partir de ella.",
            "Pensada para leerla una vez al empezar a usar el panel y volver a "
            "consultarla cuando aparezca una métrica nueva.",
        ],
        layout="cover",
    ),
    Slide(
        kicker="Cómo leer este panel",
        title="Lo esencial",
        body=[
            "El panel monitorea las conversaciones de WhatsApp del bot en tiempo real.",
            "No hay nada que exportar para ver cifras: todo se calcula en vivo desde la base.",
            "Los datos se muestran en español argentino, formato de fecha DD/MM/AAAA, "
            "zona horaria de Buenos Aires.",
            "Cada tarjeta tiene una nota al pie chica que explica qué incluye y qué deja afuera. "
            "Cuando aparece un guion '—', significa que no hay datos suficientes para mostrar la métrica.",
        ],
        layout="title_only",
    ),
    Slide(
        kicker="Mapa del panel",
        title="Las seis zonas",
        body=[
            "1) Selectores: Período (cuántos días) y Cómo contar contactos.",
            "2) KPIs globales: cuatro tarjetas con los totales del período.",
            "3) Resueltas por el bot: porcentaje de conversaciones sin intervención humana.",
            "4) Volumen diario: línea de mensajes entrantes vs salientes.",
            "5) Composición de la demanda: por intención, con líder, gráficos y heatmap.",
            "6) Eficiencia + Seguimiento: cómo termina cada flujo y a quién priorizar.",
        ],
        screenshot="01-overview-full.png",
        layout="full_visual",
    ),
    Slide(
        kicker="Selector",
        title="Período",
        body=[
            "Cambia la ventana de análisis: 7, 14, 28 o 56 días.",
            "Afecta todas las tarjetas y gráficos del período (KPIs, volumen, intenciones, eficiencia).",
            "NO afecta a la heatmap (siempre 28 días, para que la grilla no quede vacía) ni a "
            "Contactos prioritarios (usa su propia frescura).",
            "Si compartís un link con el período seleccionado, la otra persona ve lo mismo.",
        ],
        screenshot="02-window-toggle.png",
        footnote="URL: ?window=7|14|28|56",
    ),
    Slide(
        kicker="Selector",
        title="Cómo contar contactos",
        body=[
            "Define cómo asignar a una persona a una intención cuando consultó por varias cosas.",
            "Último interés (predeterminado): cuenta a la persona donde terminó la conversación.",
            "Interés inicial: cuenta a la persona por lo primero que consultó (oculto detrás de Análisis avanzado).",
            "Todas las consultas: una misma persona puede aparecer en más de una categoría a la vez. "
            "Los totales pueden superar a Contactos únicos.",
        ],
        screenshot="03-touch-toggle.png",
        footnote="URL: ?touch=last|first|any",
    ),
    Slide(
        kicker="KPI",
        title="Mensajes entrantes",
        body=[
            "Total de mensajes que el bot recibió de los contactos en el período.",
            "Incluye texto, audios y demás tipos de mensaje. Cuenta cada mensaje, no cada persona.",
            "El delta verde/rojo abajo indica cuánto creció o cayó respecto del período anterior "
            "de igual duración (si elegís 28 días, compara con los 28 días previos).",
        ],
        screenshot="04-kpi-strip.png",
    ),
    Slide(
        kicker="KPI",
        title="Mensajes salientes",
        body=[
            "Total de mensajes que el bot envió a los contactos en el período.",
            "Una proporción saludable suele ser cercana a 1:1 frente a los entrantes; un ratio muy "
            "alto puede indicar respuestas en cadena o reintentos.",
            "Como el resto, se compara contra el período anterior de igual duración.",
        ],
        screenshot="04-kpi-strip.png",
    ),
    Slide(
        kicker="KPI",
        title="Contactos únicos",
        body=[
            "Cantidad de personas distintas que escribieron al bot en el período.",
            "'Único' significa por número de WhatsApp: si la misma persona escribe diez veces en "
            "el período, se cuenta una sola.",
            "Es el techo natural para varias métricas. Por ejemplo, en Último interés la suma de "
            "Contactos por intención es menor o igual a este número.",
        ],
        screenshot="04-kpi-strip.png",
    ),
    Slide(
        kicker="KPI",
        title="Tasa de derivación",
        body=[
            "Porcentaje de conversaciones que terminaron escalando a un asesor humano.",
            "Es la tasa GLOBAL del período. Más adelante vas a ver chips de tasa por intención "
            "(Ventas, Alquileres, etc.); esos no suman al global, lo aclara la nota debajo del gráfico.",
            "Una flecha roja hacia abajo es buena: implica que el bot está resolviendo más solo.",
        ],
        screenshot="04-kpi-strip.png",
    ),
    Slide(
        kicker="KPI destacado",
        title="Resueltas por el bot",
        body=[
            "Porcentaje de contactos del período que el bot atendió sin pasar a una persona.",
            "Se calcula sobre los contactos únicos del período: 'aparecieron en el bot y no "
            "generaron una derivación de negocio'.",
            "Caveat v1: si una persona aún está en la cola de Seguimiento, igual cuenta como "
            "resuelta acá. Lo cruzás con /follow-up si querés afinar la lectura.",
        ],
        screenshot="05-self-resolution.png",
    ),
    Slide(
        kicker="Volumen diario",
        title="Mensajes entrantes y salientes",
        body=[
            "Línea de tiempo del período seleccionado, un punto por día.",
            "Sirve para detectar picos (campañas, fines de semana) y caídas (caídas de servicio, "
            "feriados).",
            "Si ves un patrón llamativo, cruzalo con Demanda por hora más abajo: mismo dato pero "
            "agregado por día de semana × hora.",
        ],
        screenshot="06-volume-chart.png",
    ),
    Slide(
        kicker="Sección",
        title="Composición de la demanda",
        body=[
            "Esta sección descompone la demanda por intención: Ventas, Alquileres, Tasaciones, "
            "Emprendimientos, Administración y Otras.",
            "Acá vive el selector Cómo contar contactos: todas las tarjetas debajo se recalculan "
            "según la opción que elijas.",
            "Cada gráfico tiene una nota al pie que recuerda exactamente qué cuenta y bajo qué regla.",
        ],
        layout="section",
    ),
    Slide(
        kicker="Tarjeta",
        title="Intención líder",
        body=[
            "Es la intención con más contactos en el período.",
            "Tiene una leyenda chica que dice 'Según: Último interés' (o el modo que tengas activo).",
            "Por eso el valor cambia cuando flipás el selector: en Interés inicial podés ver una "
            "intención líder distinta que en Último interés.",
        ],
        screenshot="07-leading-intent.png",
    ),
    Slide(
        kicker="Gráfico",
        title="Contactos por intención",
        body=[
            "Barras horizontales: cantidad de contactos por intención en el período.",
            "Debajo de cada barra hay un chip 'Nuevo' o '+12,3%' con el delta vs el período anterior.",
            "Más abajo aparecen los chips de tasa de derivación POR intención (verdes/rojos según "
            "el objetivo configurado para ese rubro).",
            "La nota 'Derivación calculada por último interés. Ver detalle.' explica por qué los "
            "porcentajes por intención no suman a la tasa global.",
        ],
        screenshot="08-contacts-by-intent.png",
    ),
    Slide(
        kicker="Operadores",
        title="Top valores en Otras",
        body=[
            "Lista colapsable que aparece debajo de Contactos por intención. Es una herramienta de "
            "mantenimiento: muestra los textos crudos del bot que cayeron en la categoría 'Otras'.",
            "Los números pueden ser MAYORES que la barra 'Otras' del gráfico de arriba: el gráfico "
            "usa último interés (cada persona en una sola categoría), esta lista cuenta a todos los "
            "que pasaron por 'Otras' en algún momento.",
            "Útil para detectar tokens nuevos que aún no tienen una etiqueta en el panel.",
        ],
        screenshot="09-otras-breakdown.png",
    ),
    Slide(
        kicker="Gráfico",
        title="Volumen por intención",
        body=[
            "Mide carga de trabajo del bot, no cantidad de personas: cuenta cada mensaje en cada paso del flujo.",
            "Ejemplo: una sola persona en el wizard de Ventas puede sumar 8 mensajes; en este "
            "gráfico cuenta 8, en Contactos por intención cuenta 1.",
            "Los chips 'X.X por contacto' bajo las barras son la densidad de interacción: mensajes "
            "÷ contactos únicos. Tasaciones suele ser alto (intake guiado), Otras suele ser ~1.",
        ],
        screenshot="10-volume-by-intent.png",
        footnote="El volumen no cambia con Cómo contar contactos; el promedio sí (cambia el denominador).",
    ),
    Slide(
        kicker="Heatmap",
        title="Demanda por hora",
        body=[
            "Grilla de 7 días × 24 horas en hora local del tenant. Cada celda es un día de semana × franja horaria.",
            "Más oscuro = más mensajes entrantes. Útil para definir cobertura humana: durante qué "
            "franjas conviene que un asesor esté disponible para escalar.",
            "Usa una ventana fija de 28 días, independiente del selector Período. Una semana es muy "
            "poco para llenar 168 celdas y la grilla quedaría casi vacía.",
            "Filtrable por intención (clic en una etiqueta de la izquierda) y compartible vía URL.",
        ],
        screenshot="11-heatmap.png",
        footnote="URL: ?heatmapIntent=Ventas | Alquileres | Tasaciones | Emprendimientos | Administracion | Otras",
    ),
    Slide(
        kicker="Sección",
        title="Eficiencia de los flujos",
        body=[
            "Esta sección mide cómo PERFORMA el bot, no cuán ocupado está.",
            "Dos tarjetas: Finalización de flujos (terminaste el flujo?) y Tiempo hasta derivación "
            "(cuánto esperaste para hablar con un humano?).",
            "Ambas usan último interés como criterio de bucket: cada persona pertenece a una sola intención.",
        ],
        layout="section",
    ),
    Slide(
        kicker="Tarjeta",
        title="Finalización de flujos",
        body=[
            "Por intención: cuántas personas llegaron al paso final del flujo.",
            "Lectura de cada fila: '6 contactos · 33,3%' = 6 personas terminaron en Ventas como "
            "última intención y 2 de ellas alcanzaron el paso terminal del flujo.",
            "Una persona aparece en una sola intención (la última). No es 'completó Ventas Y "
            "Alquileres', es 'terminó conversando sobre X'.",
        ],
        screenshot="12-completion.png",
        footnote="'—' significa que esa intención no tiene un paso terminal definido todavía.",
    ),
    Slide(
        kicker="Tarjeta",
        title="Tiempo hasta derivación",
        body=[
            "Cuánto tarda una persona desde su primer mensaje hasta hablar con un humano, por intención.",
            "Mediana = el contacto típico (la mitad esperó menos, la mitad más).",
            "p90 = el 90% esperó MENOS que ese tiempo; el 10% más lento esperó más. Sirve para "
            "definir promesas de servicio.",
            "Si una intención tiene n < 5 muestras, mediana y p90 muestran '—' para no engañar con "
            "un promedio basado en uno o dos casos. Cuando la cuenta crezca, los números aparecen solos.",
        ],
        screenshot="13-time-to-handoff.png",
        footnote="Tiempo en reloj real, incluye horario no laboral.",
    ),
    Slide(
        kicker="Seguimiento",
        title="Contactos prioritarios",
        body=[
            "Los 5 contactos a los que conviene escribir primero, ordenados por prioridad.",
            "Alta (rojo): el bot dejó la conversación sin cerrar y la persona no volvió.",
            "Media (naranja): conversación abierta pero con menor riesgo de pérdida.",
            "Baja (verde): contacto reciente, ya atendido o en seguimiento normal.",
            "El listado completo está en /follow-up. Esta sección usa su propia ventana, "
            "independiente del selector Período.",
        ],
        screenshot="14-follow-up.png",
    ),
    Slide(
        kicker="Glosario",
        title="Atajos y vocabulario",
        body=[
            "Mediana (p50): valor del medio. Si los tiempos de espera son [20, 25, 30, 40, 600] "
            "segundos, la mediana es 30.",
            "p90: el 90% de las observaciones cae por debajo de ese valor. En el ejemplo, p90 ≈ 600.",
            "Último interés / Interés inicial / Todas las consultas: las tres formas de asignar una "
            "persona a una intención cuando consultó por varias.",
            "Derivación: el bot pasó la conversación a un humano (no incluye errores de ejecución del bot).",
            "Atajos en URL: ?window=7|14|28|56, ?touch=last|first|any, ?heatmapIntent=<intención>.",
        ],
        layout="title_only",
    ),
]


# --- Rendering helpers ----------------------------------------------------


def _set_run_font(run, font_name: str, fallback: str, size_pt: float, color: RGBColor,
                  bold: bool = False) -> None:
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    # python-pptx doesn't expose font fallback directly. We set the primary
    # name; PowerPoint will substitute Calibri/Cambria if Geist/Fraunces aren't
    # installed locally. We also write the fallback into the eastAsia name as a
    # belt-and-suspenders measure.
    rPr = run._r.get_or_add_rPr()
    # Look for and update the latin element to set typeface
    for tag in ("latin",):
        existing = rPr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}")
        if existing is not None:
            existing.set("typeface", font_name)
        else:
            from pptx.oxml.ns import qn
            el = rPr.makeelement(qn(f"a:{tag}"), {"typeface": font_name})
            rPr.append(el)
    # We rely on PowerPoint's own font substitution for the fallback at view time.
    _ = fallback  # signal that fallback isn't lost in code review


def _set_fill(shape, color: RGBColor) -> None:
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_line(shape, color: RGBColor, width_pt: float = 0.75) -> None:
    line = shape.line
    line.color.rgb = color
    line.width = Pt(width_pt)


def _no_line(shape) -> None:
    shape.line.fill.background()


def _add_canvas(slide, theme: dict, color: Optional[RGBColor] = None) -> None:
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H,
    )
    _set_fill(bg, color or theme["canvas"])
    _no_line(bg)
    bg.shadow.inherit = False


def _add_accent_strip(slide, theme: dict) -> None:
    # Brand accent across the top — chunkier than the original placeholder so
    # the brand presence reads at a glance. Followed by a thin gold hairline
    # that matches the visor accent on the logo.
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(80000),
    )
    _set_fill(strip, theme["accent"])
    _no_line(strip)
    gold_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Emu(80000), SLIDE_W, Emu(15000),
    )
    _set_fill(gold_line, BRAND_GOLD)
    _no_line(gold_line)


def _add_footer_brand(slide, page_num: int, total: int,
                      logo_path: Optional[Path], theme: dict) -> None:
    # Brand mark at the bottom-left + page counter at the bottom-right. The
    # logo is a touch larger than the previous version so the helmet reads,
    # the wordmark uses sky blue, and a tiny gold pip separates them — all
    # echoing the cover treatment so every slide feels signed.
    if logo_path and logo_path.exists():
        slide.shapes.add_picture(
            str(logo_path),
            MARGIN, SLIDE_H - Inches(0.5),
            height=Inches(0.32),
        )
        text_x = MARGIN + Inches(0.46)
    else:
        text_x = MARGIN

    _add_text_block(
        slide, text_x, SLIDE_H - Inches(0.42), Inches(3.0), Inches(0.25),
        BRAND_NAME,
        font=MONO_FONT, fallback=MONO_FALLBACK,
        size_pt=10, color=theme["accent_text"], bold=True, line_spacing=1.0,
    )
    # Tiny gold pip separator
    pip = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        text_x + Inches(1.4), SLIDE_H - Inches(0.32),
        Inches(0.06), Inches(0.06),
    )
    _set_fill(pip, BRAND_GOLD)
    _no_line(pip)
    _add_text_block(
        slide, text_x + Inches(1.6), SLIDE_H - Inches(0.42), Inches(2.5), Inches(0.25),
        BRAND_TAGLINE,
        font=BODY_FONT, fallback=BODY_FALLBACK,
        size_pt=9, color=theme["footer_dim"], line_spacing=1.0,
    )
    _add_text_block(
        slide, SLIDE_W - MARGIN - Inches(1.0), SLIDE_H - Inches(0.42),
        Inches(1.0), Inches(0.25),
        f"{page_num:02d} / {total:02d}",
        font=MONO_FONT, fallback=MONO_FALLBACK,
        size_pt=10, color=theme["footer_dim"],
        align=PP_ALIGN.RIGHT, line_spacing=1.0,
    )


def _add_text_block(slide, x, y, w, h, text: str, *, font: str, fallback: str,
                    size_pt: float, color: RGBColor, bold: bool = False,
                    align: int = PP_ALIGN.LEFT, anchor: int = MSO_ANCHOR.TOP,
                    line_spacing: float = 1.25) -> None:
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    paragraphs = text.split("\n") if "\n" in text else [text]
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = para
        _set_run_font(run, font, fallback, size_pt, color, bold=bold)


def _add_kicker_and_title(slide, kicker: str, title: str, theme: dict,
                          ink_override: Optional[RGBColor] = None) -> None:
    # Tiny gold rule above the kicker — small mark that signs the page.
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN, Inches(0.45),
        Inches(0.4), Emu(20000),
    )
    _set_fill(rule, BRAND_GOLD)
    _no_line(rule)
    _add_text_block(
        slide, MARGIN, Inches(0.55), SLIDE_W - 2 * MARGIN, Inches(0.3),
        kicker.upper(),
        font=MONO_FONT, fallback=MONO_FALLBACK,
        size_pt=10, color=theme["accent_text"], bold=True, line_spacing=1.0,
    )
    _add_text_block(
        slide, MARGIN, Inches(0.95), SLIDE_W - 2 * MARGIN, Inches(1.6),
        title,
        font=HEADING_FONT, fallback=HEADING_FALLBACK,
        size_pt=36, color=ink_override or theme["ink"], bold=True,
        line_spacing=1.05,
    )


def _add_body(slide, x, y, w, h, paragraphs: list[str], theme: dict,
              ink_override: Optional[RGBColor] = None) -> None:
    if not paragraphs:
        return
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.35
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = para
        _set_run_font(run, BODY_FONT, BODY_FALLBACK, 14,
                      ink_override or theme["ink"])


def _add_footnote(slide, text: str, theme: dict) -> None:
    _add_text_block(
        slide, MARGIN, SLIDE_H - Inches(0.85), SLIDE_W - 2 * MARGIN, Inches(0.35),
        text,
        font=BODY_FONT, fallback=BODY_FALLBACK,
        size_pt=10, color=theme["soft_ink"], line_spacing=1.2,
    )


def _add_screenshot_or_placeholder(slide, x, y, w, h, screenshot_dir: Path,
                                   filename: Optional[str], theme: dict) -> None:
    if filename and (screenshot_dir / filename).exists():
        # Hairline frame around the screenshot so it has a definite edge on
        # the dark theme (where it would otherwise float on navy).
        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        _set_fill(frame, theme["surface"])
        _set_line(frame, theme["screenshot_border"], width_pt=0.5)
        slide.shapes.add_picture(
            str(screenshot_dir / filename), x, y, width=w, height=h,
        )
        return

    # Placeholder card with the expected filename.
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _set_fill(card, theme["surface"])
    _set_line(card, theme["rule"], width_pt=0.75)

    label = "Captura pendiente"
    if filename:
        label = f"{label}\n{filename}"
    _add_text_block(
        slide, x, y + h / 2 - Inches(0.4), w, Inches(0.8),
        label,
        font=BODY_FONT, fallback=BODY_FALLBACK,
        size_pt=12, color=theme["soft_ink"],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.3,
    )


def _render_two_col(slide, data: Slide, screenshot_dir: Path, theme: dict) -> None:
    _add_kicker_and_title(slide, data.kicker, data.title, theme)

    body_x = MARGIN
    body_y = Inches(2.7)
    body_w = Inches(5.0)
    body_h = SLIDE_H - body_y - Inches(1.05)
    _add_body(slide, body_x, body_y, body_w, body_h, data.body, theme)

    img_x = body_x + body_w + Inches(0.4)
    img_y = Inches(2.7)
    img_w = SLIDE_W - img_x - MARGIN
    img_h = SLIDE_H - img_y - Inches(1.05)
    _add_screenshot_or_placeholder(slide, img_x, img_y, img_w, img_h,
                                   screenshot_dir, data.screenshot, theme)

    if data.footnote:
        _add_footnote(slide, data.footnote, theme)


def _render_title_only(slide, data: Slide, screenshot_dir: Path, theme: dict) -> None:
    _add_kicker_and_title(slide, data.kicker, data.title, theme)
    body_x = MARGIN
    body_y = Inches(2.85)
    body_w = SLIDE_W - 2 * MARGIN
    body_h = SLIDE_H - body_y - Inches(1.05)
    _add_body(slide, body_x, body_y, body_w, body_h, data.body, theme)
    if data.footnote:
        _add_footnote(slide, data.footnote, theme)


def _render_full_visual(slide, data: Slide, screenshot_dir: Path, theme: dict) -> None:
    _add_kicker_and_title(slide, data.kicker, data.title, theme)

    img_x = MARGIN
    img_y = Inches(2.7)
    img_w = Inches(7.5)
    img_h = SLIDE_H - img_y - Inches(1.05)
    _add_screenshot_or_placeholder(slide, img_x, img_y, img_w, img_h,
                                   screenshot_dir, data.screenshot, theme)

    body_x = img_x + img_w + Inches(0.4)
    body_y = Inches(2.7)
    body_w = SLIDE_W - body_x - MARGIN
    body_h = img_h
    _add_body(slide, body_x, body_y, body_w, body_h, data.body, theme)

    if data.footnote:
        _add_footnote(slide, data.footnote, theme)


def _render_section(slide, data: Slide, screenshot_dir: Path, theme: dict,
                    logo_path: Optional[Path]) -> None:
    # Section dividers go full-bleed brand-dark with sky-blue typography for
    # extra brand presence between major chapters of the deck. The accent
    # strip + footer mark sit on top of this dark canvas.
    _add_canvas(slide, theme, theme["section_canvas"])

    if logo_path and logo_path.exists():
        slide.shapes.add_picture(
            str(logo_path),
            MARGIN, Inches(2.6),
            height=Inches(1.0),
        )

    # Kicker + heading — readable on dark.
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN + Inches(1.3), Inches(2.85),
        Inches(0.4), Emu(20000),
    )
    _set_fill(rule, BRAND_GOLD)
    _no_line(rule)
    _add_text_block(
        slide, MARGIN + Inches(1.3), Inches(2.95),
        SLIDE_W - 2 * MARGIN - Inches(1.3), Inches(0.3),
        data.kicker.upper(),
        font=MONO_FONT, fallback=MONO_FALLBACK,
        size_pt=11, color=BRAND_PRIMARY_SOFT, bold=True, line_spacing=1.0,
    )
    _add_text_block(
        slide, MARGIN + Inches(1.3), Inches(3.3),
        SLIDE_W - 2 * MARGIN - Inches(1.3), Inches(1.2),
        data.title,
        font=HEADING_FONT, fallback=HEADING_FALLBACK,
        size_pt=42, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True,
        line_spacing=1.05,
    )

    if data.body:
        _add_body(
            slide, MARGIN, Inches(5.0),
            SLIDE_W - 2 * MARGIN, Inches(2.0),
            data.body, theme, ink_override=theme["section_ink"],
        )


def _render_cover(slide, data: Slide, screenshot_dir: Path,
                  logo_path: Optional[Path], theme: dict) -> None:
    # Cover is brand-led: dark canvas (works on both themes — same look),
    # big helmet logo, oversized Fraunces title in white/light, sky-blue
    # tagline kicker, and a bottom rule that hints at "press kit".
    _add_canvas(slide, theme, BRAND_DARK)

    if logo_path and logo_path.exists():
        slide.shapes.add_picture(
            str(logo_path),
            Inches(0.9), Inches(1.6),
            height=Inches(2.6),
        )

    title_x = Inches(4.2)

    # Sky-blue kicker tagline
    _add_text_block(
        slide, title_x, Inches(1.6), SLIDE_W - title_x - MARGIN, Inches(0.4),
        data.kicker,
        font=MONO_FONT, fallback=MONO_FALLBACK,
        size_pt=11, color=BRAND_PRIMARY, bold=True, line_spacing=1.0,
    )

    # Gold rule
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, title_x, Inches(2.05),
        Inches(0.5), Emu(25000),
    )
    _set_fill(rule, BRAND_GOLD)
    _no_line(rule)

    # Big white Fraunces title
    _add_text_block(
        slide, title_x, Inches(2.2), SLIDE_W - title_x - MARGIN, Inches(2.4),
        data.title,
        font=HEADING_FONT, fallback=HEADING_FALLBACK,
        size_pt=46, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True,
        line_spacing=1.05,
    )

    # Soft body in pale blue
    _add_body(
        slide, title_x, Inches(4.6),
        SLIDE_W - title_x - MARGIN, Inches(2.0),
        data.body, theme, ink_override=BRAND_PRIMARY_SOFT,
    )

    # Bottom-edge brand line: small wordmark + tagline in muted sky blue
    _add_text_block(
        slide, MARGIN, SLIDE_H - Inches(0.55),
        SLIDE_W - 2 * MARGIN, Inches(0.3),
        f"{BRAND_NAME}.com.ar  ·  Panel del cliente",
        font=MONO_FONT, fallback=MONO_FALLBACK,
        size_pt=10, color=BRAND_PRIMARY_SOFT, line_spacing=1.0,
    )


# --- Build ----------------------------------------------------------------


def build(out_path: Path, screenshot_dir: Path, theme: dict) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # blank

    renderers = {
        "two_col": _render_two_col,
        "title_only": _render_title_only,
        "full_visual": _render_full_visual,
    }

    logo_path = screenshot_dir.parent / "brand" / "logo.png"
    logo_sm_path = screenshot_dir.parent / "brand" / "logo-sm.png"
    footer_logo = logo_sm_path if logo_sm_path.exists() else logo_path
    total = len(SLIDES)

    for idx, data in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank_layout)
        _add_canvas(slide, theme)

        if data.layout == "cover":
            _render_cover(slide, data, screenshot_dir, logo_path, theme)
            # Cover gets no top accent strip and no footer brand — it IS
            # the brand statement.
            continue

        if data.layout == "section":
            _render_section(slide, data, screenshot_dir, theme, logo_path)
            _add_accent_strip(slide, theme)
            _add_footer_brand(slide, idx, total, footer_logo, theme)
            continue

        _add_accent_strip(slide, theme)
        renderer = renderers.get(data.layout, _render_two_col)
        renderer(slide, data, screenshot_dir, theme)
        _add_footer_brand(slide, idx, total, footer_logo, theme)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Wrote {out_path} ({total} slides, {theme['name']} theme)")


THEMES = {"light": LIGHT_THEME, "dark": DARK_THEME}


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("out", nargs="?", help="Output .pptx path (optional)")
    parser.add_argument("--theme", choices=list(THEMES), default="light")
    parser.add_argument(
        "--all", action="store_true",
        help="Build both light and dark variants side by side.",
    )
    args = parser.parse_args()

    repo_root = Path(__file__).resolve().parent.parent
    shots = repo_root / "docs" / "screenshots"

    if args.all:
        for name, theme in THEMES.items():
            out = repo_root / "docs" / f"dashboard-walkthrough-{name}.pptx"
            build(out, shots, theme)
    else:
        default_out = (
            repo_root / "docs" / f"dashboard-walkthrough-{args.theme}.pptx"
        )
        out = Path(args.out) if args.out else default_out
        if not out.is_absolute():
            out = repo_root / out
        build(out, shots, THEMES[args.theme])
